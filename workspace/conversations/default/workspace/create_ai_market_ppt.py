#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2024年中国AI市场分析PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 76, 129)  # 深蓝色
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(200, 220, 240)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, layout_type='bullet'):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 76, 129)
    
    # 添加装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.15), Inches(9), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(15, 76, 129)
    line.line.color.rgb = RGBColor(15, 76, 129)
    
    # 添加内容
    if layout_type == 'bullet':
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(64, 64, 64)
            p.level = 0
            p.space_before = Pt(12)
            p.line_spacing = 1.3
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 76, 129)
    
    # 添加装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.15), Inches(9), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(15, 76, 129)
    line.line.color.rgb = RGBColor(15, 76, 129)
    
    # 添加表格
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(0.8), Inches(1.8),
        Inches(8.4), Inches(4.5)
    ).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 76, 129)
        
        # 设置文字格式
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # 填充数据
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            
            # 设置文字格式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
                paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.LEFT
    
    return slide

def create_ai_market_ppt():
    """创建2024年中国AI市场分析PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 封面页
    create_title_slide(
        prs,
        "2024年中国AI市场分析",
        "市场规模 · 竞争格局 · 技术趋势 · 投资机会"
    )
    
    # 2. 目录页
    add_content_slide(
        prs,
        "目 录",
        [
            "01  市场规模与增长趋势",
            "02  主要玩家与竞争格局",
            "03  关键技术趋势",
            "04  投资机会分析"
        ]
    )
    
    # 3. 市场规模与增长趋势 - 总体规模
    add_content_slide(
        prs,
        "01 市场规模与增长趋势",
        [
            "• 整体市场规模：2024年中国AI产业规模突破7000亿元，达到7470亿元",
            "  数据来源：艾媒咨询《2024-2025年中国人工智能行业发展趋势研究报告》",
            "",
            "• 增长态势：同比增长41%，连续多年保持20%以上高速增长",
            "  数据来源：中国互联网络信息中心《中国互联网络发展状况统计报告》",
            "",
            "• 市场预测：2025-2029年将保持32.1%的年均复合增长率",
            "  预计2029年突破1万亿规模",
            "  数据来源：艾瑞咨询《2024年中国人工智能产业研究报告》"
        ]
    )
    
    # 4. 市场规模 - 细分领域
    add_table_slide(
        prs,
        "市场细分领域规模",
        ["细分领域", "2024年市场规模", "增长率", "数据来源"],
        [
            ["AI大模型", "294.16亿元", "预计2026年破700亿", "艾媒咨询"],
            ["AI芯片", "1447亿元", "79.9%(五年复合)", "前瞻产业研究院"],
            ["AI公有云服务", "195.9亿元", "55.3%", "IDC"],
            ["计算机视觉", "81.0亿元", "33.7%", "IDC"],
            ["自然语言处理", "22.2亿元", "51.1%", "IDC"],
            ["AIGC市场", "436亿元", "预计2030年破1.14万亿", "Statista"]
        ]
    )
    
    # 5. 主要玩家与竞争格局 - 头部企业
    add_content_slide(
        prs,
        "02 主要玩家与竞争格局",
        [
            "• 互联网科技巨头（BAT+）",
            "  - 百度：文心一言大模型，百度智能云市场第二",
            "  - 阿里：通义千问，云计算基础设施领先",
            "  - 腾讯：混元大模型，腾讯云计算机视觉市场第一",
            "  - 华为：盘古大模型，全栈技术布局",
            "",
            "• AI原生企业",
            "  - 商汤科技：AI大装置SenseCore，传统计算机视觉领先",
            "  - 月之暗面(Moonshot)：Kimi大模型，2024年获超10亿美元融资",
            "  - DeepSeek：开源低成本大模型，2025年初引发行业震动"
        ]
    )
    
    # 6. 竞争格局 - 区域分布
    add_content_slide(
        prs,
        "区域竞争格局",
        [
            "• 三大核心极点：北京、广东(深圳)、上海",
            "  - 北京：企业扶持、技术探索、产业落地综合排名第一",
            "  - 广东：AI创新能力80.5分位居榜首",
            "  - 浙江：2024年首次跻身引领者梯队，区域潜力第一",
            "  数据来源：中国电子信息产业发展研究院《中国人工智能区域竞争力研究报告》",
            "",
            "• 产业集群特征：",
            "  - 以华为、腾讯、京东、阿里及三大运营商为核心节点",
            "  - 形成复杂的技术合作生态网络",
            "  - 呈现'极化'与'扩散'并存的发展态势"
        ]
    )
    
    # 7. 关键技术趋势
    add_content_slide(
        prs,
        "03 关键技术趋势",
        [
            "• 大模型技术持续演进",
            "  - Transformer架构主导，Scaling Law仍有提升空间",
            "  - 多模态融合：文本+图像+视频+音频全模态理解与生成",
            "  - 推理优化：后训练思维链(CoT)，从预训练转向强化学习",
            "",
            "• 生成式AI(AIGC)应用爆发",
            "  - 营销场景应用占比57.3%，办公49%，教育40.8%",
            "  - 346款生成式AI服务完成备案(截至2024年3月)",
            "  - 80.9%用户使用AI产品回答问题，满意度高",
            "  数据来源：创业邦《2024年AIGC创新应用洞察报告》"
        ]
    )
    
    # 8. 技术趋势 - AI Agent与端侧AI
    add_content_slide(
        prs,
        "技术趋势：AI Agent与端侧AI",
        [
            "• AI Agent（智能体）成为新方向",
            "  - 具备自主感知、决策和执行能力",
            "  - 代表产品：智谱AutoGLM、蝴蝶效应Manus",
            "  - 应用场景：企业研究、旅行规划、课程设计等",
            "",
            "• 端侧AI寻求突破",
            "  - 端侧大模型部署在手机、PC等终端",
            "  - 优势：本地化运行、隐私保护强",
            "  - 应用：手机AI助手、智能汽车等",
            "",
            "• 工具生态完善",
            "  - 分布式AI框架：DeepSpeed、Megatron、Colossal-AI",
            "  - LLMOps平台、AI一体机加速大模型部署"
        ]
    )
    
    # 9. 投资机会分析 - 投资规模
    add_content_slide(
        prs,
        "04 投资机会分析",
        [
            "• 投资规模持续增长",
            "  - 2024年人工智能行业融资超1000亿元",
            "  - 2024年上半年AIGC领域142次融资事件，披露金额142.8亿元",
            "  - 50%的AI公司在成立三年内获得投资",
            "  数据来源：IT桔子、睿兽分析",
            "",
            "• 国家级基金支持",
            "  - 国家人工智能产业投资基金成立，出资额600.6亿元",
            "  - 国务院：加大AI领域金融和财政支持，发展长期资本、耐心资本",
            "  数据来源：天眼查、国务院政策文件"
        ]
    )
    
    # 10. 投资热点赛道
    add_table_slide(
        prs,
        "投资热点赛道",
        ["投资方向", "热度", "代表企业/案例"],
        [
            ["通用大模型", "极高", "月之暗面(10亿美元)、DeepSeek"],
            ["AI芯片", "高", "11起融资事件(2024上半年)"],
            ["算力基础设施", "高", "智算中心、AI服务器"],
            ["AIGC应用层", "最热", "营销、办公、教育场景"],
            ["AI基础数据服务", "中高", "6起融资事件(2024上半年)"],
            ["垂直行业应用", "中高", "金融、医疗、工业制造"]
        ]
    )
    
    # 11. 投资机会与建议
    add_content_slide(
        prs,
        "投资机会与建议",
        [
            "• 应用层投资机会最大",
            "  - 应用层客户留存率表现更突出",
            "  - 知识管理场景最受企业青睐(52%)",
            "  - 营销、办公、教育为热门应用方向",
            "",
            "• 区域布局建议",
            "  - 超一线城市(北京、深圳、上海、杭州)产业生态最完善",
            "  - 76.4%企业认为产业生态是区域拓展首要考量因素",
            "  - 长三角、大湾区、珠三角、京津冀为关键区域",
            "",
            "• 出海机会",
            "  - 近六成AIGC企业已有海外布局",
            "  - 中国AI专利占全球60%，专利申请量全球第一(38.58%)"
        ]
    )
    
    # 12. 总结页
    add_content_slide(
        prs,
        "核心观点总结",
        [
            "✓ 市场规模：2024年破7000亿，高速增长态势持续",
            "",
            "✓ 竞争格局：BAT+华为领跑，AI原生企业快速崛起，区域集中于京粤沪",
            "",
            "✓ 技术趋势：大模型多模态演进，AIGC应用爆发，AI Agent成为新方向",
            "",
            "✓ 投资机会：应用层最热，通用大模型、AI芯片、垂直场景为重点赛道",
            "",
            "✓ 未来展望：2029年市场规模将突破1万亿，AI+产业融合深化"
        ]
    )
    
    # 13. 致谢页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 76, 129)
    
    thank_you = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    thank_frame = thank_you.text_frame
    thank_frame.text = "感谢观看"
    thank_para = thank_frame.paragraphs[0]
    thank_para.font.size = Pt(48)
    thank_para.font.bold = True
    thank_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_para.alignment = PP_ALIGN.CENTER
    
    source = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
    source_frame = source.text_frame
    source_frame.text = "数据来源：艾瑞咨询、艾媒咨询、IDC、中国信通院等权威机构"
    source_para = source_frame.paragraphs[0]
    source_para.font.size = Pt(14)
    source_para.font.color.rgb = RGBColor(200, 220, 240)
    source_para.alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_path = "2024年中国AI市场分析报告.pptx"
    prs.save(output_path)
    print(f"✅ PPT创建成功！文件保存在：{os.path.abspath(output_path)}")
    return output_path

if __name__ == "__main__":
    create_ai_market_ppt()
