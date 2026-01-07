#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2024年中国AI市场分析PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 73, 125)  # 深蓝色
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(10), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(31, 73, 125)
    title_shape.line.color.rgb = RGBColor(31, 73, 125)
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.margin_top = Inches(0.1)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_bullet_points(slide, left, top, width, height, items, font_size=18):
    """添加项目符号列表"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
    
    return textbox

def add_chart_bar(slide, left, top, width, height, categories, values, chart_title=""):
    """添加柱状图"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('数值', values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart
    
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    return chart

def add_data_box(slide, left, top, width, height, title, value, unit="", bg_color=None):
    """添加数据展示框"""
    if bg_color is None:
        bg_color = RGBColor(68, 114, 196)
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = bg_color
    
    text_frame = box.text_frame
    text_frame.clear()
    
    # 标题
    p1 = text_frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(10)
    
    # 数值
    p2 = text_frame.add_paragraph()
    p2.text = value + unit
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    return box

def create_ppt():
    """生成完整PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # ============ 第1页：封面 ============
    create_title_slide(
        prs,
        "2024年中国AI市场分析",
        "市场规模 · 竞争格局 · 技术趋势 · 投资机会"
    )
    
    # ============ 第2页：目录 ============
    slide = create_content_slide(prs, "目录")
    items = [
        "01  市场规模和增长趋势",
        "02  主要玩家和竞争格局",
        "03  关键技术趋势",
        "04  投资机会分析"
    ]
    add_bullet_points(slide, 2, 1.5, 6, 3, items, font_size=24)
    
    # ============ 第3页：市场规模概览 ============
    slide = create_content_slide(prs, "01 | 市场规模和增长趋势")
    
    # 关键数据展示
    add_data_box(slide, 0.5, 1.2, 2.5, 1.2, "2024年市场规模", "7470", "亿元", RGBColor(68, 114, 196))
    add_data_box(slide, 3.2, 1.2, 2.5, 1.2, "同比增长", "41", "%", RGBColor(237, 125, 49))
    add_data_box(slide, 5.9, 1.2, 2.5, 1.2, "全球占比", "20.9", "%", RGBColor(112, 173, 71))
    
    # 要点说明
    points = [
        "✓ 2024年中国AI产业规模突破7000亿元，连续多年保持20%以上增长",
        "✓ 大模型产业规模达216亿元，年均复合增长率116%",
        "✓ AIGC市场规模436亿元，预计2030年达1.14万亿元",
        "✓ AI公有云服务市场195.9亿元，同比增长55.3%"
    ]
    add_bullet_points(slide, 1, 2.8, 8, 2.5, points, font_size=16)
    
    # ============ 第4页：细分市场 ============
    slide = create_content_slide(prs, "01 | 细分市场结构")
    
    # 细分领域数据
    categories = ['计算机视觉', '机器学习', '自然语言', 'MaaS服务']
    values = [81, 7.1, 195.9, 216]
    add_chart_bar(slide, 0.8, 1.5, 8, 3.5, categories, values, "主要细分市场规模(亿元)")
    
    # ============ 第5页:竞争格局 ============
    slide = create_content_slide(prs, "02 | 主要玩家和竞争格局")
    
    # 左侧：第一梯队企业
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "第一梯队：四大头部企业"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 73, 125)
    
    companies = [
        "🔹 字节跳动（豆包）：46.4%市场份额，日均调用12.7万亿Token",
        "🔹 阿里巴巴（通义千问）：19.3%市场份额，夸克月活过亿",
        "🔹 百度（文心一言）：日均调用16.5亿次，AI收入增长300%",
        "🔹 腾讯（混元）：600+内部场景落地，元宝月活2636万"
    ]
    add_bullet_points(slide, 0.5, 1.7, 4.5, 3, companies, font_size=15)
    
    # 右侧：竞争态势
    title_box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(0.4))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "竞争态势特征"
    p2 = title_frame2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(31, 73, 125)
    
    trends = [
        "💡 价格战激烈：主力模型降价97%，轻量模型免费",
        "💡 生态为王：从模型能力转向应用生态建设",
        "💡 场景深耕：金融、医疗、政务等垂直领域渗透",
        "💡 技术分化：开源vs闭源，通用vs垂直"
    ]
    add_bullet_points(slide, 5.2, 1.7, 4.3, 3, trends, font_size=15)
    
    # ============ 第6页：技术趋势 ============
    slide = create_content_slide(prs, "03 | 关键技术趋势")
    
    # 三栏布局
    trends_data = [
        {
            "title": "🚀 大模型进化",
            "items": [
                "参数规模突破万亿级",
                "多模态融合成主流",
                "MoE架构降低成本",
                "推理速度大幅提升"
            ],
            "left": 0.5
        },
        {
            "title": "🎨 AIGC应用爆发",
            "items": [
                "文本/图像/视频生成",
                "营销场景占比57%",
                "办公教育需求旺盛",
                "知识管理成企业首选"
            ],
            "left": 3.7
        },
        {
            "title": "⚡ 基础设施升级",
            "items": [
                "智算中心加速建设",
                "东数西算工程推进",
                "AI芯片技术突破",
                "算力成本持续下降"
            ],
            "left": 6.9
        }
    ]
    
    for trend in trends_data:
        # 标题
        title_box = slide.shapes.add_textbox(Inches(trend["left"]), Inches(1.2), Inches(2.8), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = trend["title"]
        p = title_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 73, 125)
        
        # 内容
        add_bullet_points(slide, trend["left"], 1.8, 2.8, 3, trend["items"], font_size=14)
    
    # ============ 第7页：投资机会 ============
    slide = create_content_slide(prs, "04 | 投资机会分析")
    
    # 关键数据
    add_data_box(slide, 0.5, 1.2, 2.3, 1.1, "2024年融资规模", "1000", "亿元+", RGBColor(237, 125, 49))
    add_data_box(slide, 3, 1.2, 2.3, 1.1, "融资事件", "142", "起", RGBColor(112, 173, 71))
    add_data_box(slide, 5.5, 1.2, 2.3, 1.1, "国家AI基金", "600.6", "亿元", RGBColor(68, 114, 196))
    
    # 投资热点
    hotspots = [
        "🔥 热门赛道：通用大模型、AI芯片、文娱应用领跑",
        "🔥 地域分布：北京、深圳、上海、杭州聚集度最高",
        "🔥 应用层机会：营销、办公、教育场景客户留存率突出",
        "🔥 政策支持：多地出台大模型产业发展措施，智算中心投资拉动产业增长3-4倍",
        "🔥 出海趋势：近60%企业已有海外布局，产业生态成关键考量因素"
    ]
    add_bullet_points(slide, 0.8, 2.6, 8.5, 2.8, hotspots, font_size=15)
    
    # ============ 第8页：总结与展望 ============
    slide = create_content_slide(prs, "总结与展望")
    
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(3.5)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
    summary_box.line.color.rgb = RGBColor(68, 114, 196)
    summary_box.line.width = Pt(2)
    
    text_frame = summary_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.3)
    
    summaries = [
        "📊 市场现状：2024年中国AI市场规模7470亿元，同比增长41%，全球占比超20%",
        "",
        "🏆 竞争格局：字节、阿里、百度、腾讯四大头部企业引领，价格战推动应用普及",
        "",
        "🔬 技术演进：大模型向万亿参数进化，多模态融合，AIGC应用全面爆发",
        "",
        "💰 投资前景：2024年融资超千亿，国家级基金入场，应用层和基础设施成投资重点",
        "",
        "🚀 未来预测：预计2025年市场规模突破万亿，AI将深度渗透各行业，成为新质生产力"
    ]
    
    for i, text in enumerate(summaries):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = text
        if text:
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(8)
    
    # ============ 第9页：数据来源 ============
    slide = create_content_slide(prs, "数据来源")
    
    sources = [
        "📚 中国信息通信研究院：《人工智能发展报告(2024年)》",
        "📚 中国电子信息产业发展研究院（赛迪研究院）：《中国人工智能区域竞争力研究报告》",
        "📚 艾媒咨询：《2024-2025年中国人工智能行业发展趋势研究报告》",
        "📚 IDC：《中国模型即服务(MaaS)及AI大模型解决方案市场追踪，2024H2》",
        "📚 前瞻产业研究院：《2024年中国生成式AI行业全景图谱》",
        "📚 创业邦研究中心：《2024年AIGC创新应用洞察报告》",
        "📚 IT桔子、睿兽分析：AI投融资数据",
        "📚 百度、阿里、腾讯、字节跳动等企业公开财报及发布会资料"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 3.5, sources, font_size=15)
    
    # 保存
    output_path = '/Users/liuyi/Documents/langchain/CoT_agent/mvp/zenflux_agent/workspace/conversations/default/workspace/2024年中国AI市场分析.pptx'
    prs.save(output_path)
    print(f"✅ PPT已成功生成：{output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
