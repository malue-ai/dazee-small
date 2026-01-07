from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """添加封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色（渐变效果用纯色替代）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "2024年度深度研究报告"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 126, 234)
    
    # 添加下划线
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(102, 126, 234)
    line.line.width = Pt(3)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.level = 0

# 1. 封面
add_title_slide(
    "2024年中国AI市场分析",
    "市场规模 · 竞争格局 · 技术趋势 · 投资机会"
)

# 2. 目录
add_content_slide("目录", [
    "01  市场规模和增长趋势",
    "      整体市场规模、细分领域、增长动力",
    "",
    "02  主要玩家和竞争格局",
    "      头部企业、市场份额、竞争态势",
    "",
    "03  关键技术趋势",
    "      大模型、AIGC、AI Agent、多模态",
    "",
    "04  投资机会",
    "      融资数据、热门赛道、投资建议"
])

# 3. 市场规模
add_content_slide("市场规模和增长趋势", [
    "🚀 整体市场规模",
    "• 2024年中国AI行业市场规模：7,470亿元",
    "• 同比增长率：41%",
    "• AI公有云服务市场规模：195.9亿元（+55.3%）",
    "",
    "📈 细分市场表现",
    "• AI大模型：294.16亿元",
    "• AI芯片：1,447亿元",
    "• 计算机视觉：81亿元（+33.7%）",
    "• 自然语言处理：22.2亿元（+51.1%）",
    "",
    "🎯 增长预测",
    "• 2026年大模型市场将突破700亿元",
    "• 2035年生成式AI将超30万亿元",
    "",
    "数据来源：艾媒咨询、IDC、前瞻产业研究院"
])

# 4. 市场驱动力
add_content_slide("市场增长驱动力", [
    "🏛️ 政策支持",
    "• \"人工智能+\"行动全面实施",
    "• \"十四五\"规划重点布局AI产业",
    "• \"东数西算\"工程推进算力建设",
    "",
    "💻 算力基础",
    "• 全国30+城市建设智算中心",
    "• 云计算市场达3097.3亿元",
    "• AI芯片年复合增长率79.9%",
    "",
    "🔬 技术突破",
    "• 大模型技术快速迭代",
    "• 多模态AI能力持续提升",
    "• AI Agent自主性增强",
    "",
    "📱 应用普及",
    "• 金融、医疗等行业渗透率超50%",
    "• 企业数字化转型需求旺盛",
    "",
    "数据来源：中国信通院《人工智能发展报告(2024年)》"
])

# 5. 竞争格局
add_content_slide("主要玩家和竞争格局", [
    "🔵 百度 - 文心大模型系列",
    "• 文心一言用户超1亿  • 搜索+AI深度融合",
    "",
    "🟠 阿里巴巴 - 通义千问系列",
    "• 电商+AI场景优势  • 云服务生态整合",
    "",
    "🟢 腾讯 - 混元大模型",
    "• 社交+内容AI赋能  • 计算机视觉市场第一",
    "",
    "🔴 华为 - 盘古大模型",
    "• To B行业大模型专家  • 算力+芯片全栈布局",
    "",
    "🌟 新兴力量",
    "字节跳动（豆包）、科大讯飞（星火）、商汤科技、",
    "旷视科技、月之暗面（Kimi）、百川智能、智谱AI等",
    "",
    "数据来源：艾瑞咨询《2024年中国人工智能产业研究报告》"
])

# 6. 技术趋势
add_content_slide("关键技术趋势", [
    "🤖 大模型（LLM）",
    "• 参数规模持续扩大：千亿级成主流",
    "• 多模态融合：文本、图像、音频统一",
    "• 垂直领域模型：金融、医疗专用大模型",
    "• 端侧大模型：手机、PC本地化运行",
    "",
    "🎨 生成式AI（AIGC）",
    "• 内容创作：文案、设计、视频生成",
    "• 代码生成：提升开发效率",
    "• 2024市场规模14.4万亿元",
    "",
    "🦾 AI Agent（智能体）",
    "• 自主决策能力：无需人工干预",
    "• 任务分解执行：复杂任务自动完成",
    "• 2025进入Agent时代",
    "",
    "数据来源：中国软件评测中心、腾讯研究院"
])

# 7. 投资机会
add_content_slide("投资机会", [
    "📊 投融资数据",
    "• 2024年AI行业融资总额：1,000亿元+",
    "• 50%的AI公司在成立3年内获得投资",
    "",
    "🔥 热门投资赛道",
    "• 大模型基础设施：算力、芯片、训练平台",
    "• 垂直行业应用：金融、医疗、教育",
    "• AIGC工具：内容生成、设计工具",
    "• AI Agent平台：智能体开发与运营",
    "",
    "🎯 高渗透率行业",
    "• 金融：渗透率>50%，风控+投顾",
    "• 政府：渗透率>50%，智慧城市",
    "• 教育：渗透率>50%，个性化学习",
    "• 医疗：影像诊断+新药研发",
    "",
    "数据来源：清科研究中心、IT桔子、投中网"
])

# 8. 总结
add_content_slide("核心结论", [
    "🌟 市场展望",
    "",
    "✅ 爆发式增长",
    "2024年市场规模7,470亿元，未来3年CAGR超40%",
    "",
    "✅ 政策红利",
    "\"人工智能+\"行动持续推进，算力基础设施加速建设",
    "",
    "✅ 技术成熟",
    "大模型进入应用落地期，2025年开启AI Agent时代",
    "",
    "✅ 竞争激烈",
    "BAT+华为领跑，新兴企业快速崛起",
    "",
    "✅ 投资活跃",
    "年度融资超千亿，垂直应用成投资热点",
    "",
    "━━━━━━━━━━━━━━━━━━━━━━",
    "🚀 2024-2026是中国AI产业的黄金发展期"
])

# 保存文件
output_path = '/Users/liuyi/Documents/langchain/CoT_agent/mvp/zenflux_agent/workspace/conversations/default/workspace/2024年中国AI市场分析.pptx'
prs.save(output_path)
print(f"✅ PPT已成功创建：{output_path}")
print(f"📊 共生成 {len(prs.slides)} 页幻灯片")
