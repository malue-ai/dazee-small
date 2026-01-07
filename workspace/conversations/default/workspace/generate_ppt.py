#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定义颜色方案
COLOR_TITLE = RGBColor(31, 78, 121)  # 深蓝色
COLOR_ACCENT = RGBColor(0, 176, 240)  # 浅蓝色
COLOR_TEXT = RGBColor(64, 64, 64)  # 深灰色

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加背景色
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 248, 255)
    background.line.fill.background()
    
    # 添加主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLOR_ACCENT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE
    
    # 添加下划线
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, content in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = content
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# 第1页：封面
add_title_slide(prs, "2024年中国AI市场分析", "市场规模 · 竞争格局 · 技术趋势 · 投资机会")

# 第2页：目录
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "目录"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_TITLE

content_items = [
    "01  市场规模和增长趋势",
    "02  主要玩家和竞争格局", 
    "03  关键技术趋势",
    "04  投资机会分析"
]

y_pos = 2.0
for i, item in enumerate(content_items):
    content_box = slide2.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.6))
    text_frame = content_box.text_frame
    text_frame.text = item
    p = text_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_TEXT
    y_pos += 0.9

# 第3页：市场规模概览
add_content_slide(prs, "一、市场规模和增长趋势", [
    "整体市场规模",
    "• 2024年中国AI产业规模突破7000亿元",
    "• 连续多年保持20%以上的高速增长率",
    "• AI基础设施市场规模约1513.9亿元（2023年）",
    "",
    "细分市场表现",
    "• AI大模型市场：294.16亿元（2024年），预计2026年突破700亿元",
    "• AI芯片市场：预计2024年达到1447亿元，GPU芯片占比85%",
    "• AIGC应用市场：436亿元，预计2030年达到1.14万亿元",
    "",
    "数据来源：中国信通院、艾媒咨询、IDC"
])

# 第4页：市场增长驱动因素
add_content_slide(prs, "市场增长驱动因素", [
    "政策支持",
    "• AI+战略写入政府工作报告",
    "• 国家人工智能产业投资基金成立，出资额600.6亿元",
    "• 多个城市出台AI发展鼓励政策",
    "",
    "技术突破",
    "• 大模型技术快速迭代，参数规模持续扩大",
    "• 推理成本大幅下降，性能显著提升",
    "• 多模态AI、具身智能等前沿技术突破",
    "",
    "应用场景拓展",
    "• 从单点应用向生成式AI全面升级",
    "• 深度渗透金融、医疗、教育、制造等行业"
])

# 第5页：主要玩家和竞争格局
add_content_slide(prs, "二、主要玩家和竞争格局", [
    "AI云市场份额（2025年H1）",
    "• 阿里云：35.8%（市场领先，通义千问）",
    "• 字节火山引擎：14.8%（豆包大模型）",
    "• 华为云：13.1%（盘古大模型）",
    "• 腾讯云：7.0%（混元大模型）",
    "• 百度智能云：6.1%（文心一言）",
    "",
    "大模型应用市场（2024年）",
    "• 百度：市场份额19.9%，位居第一",
    "• 商汤科技：16.0%，紧随其后",
    "• 智谱AI：初创企业中的佼佼者",
    "",
    "数据来源：IDC、Canalys"
])

# 第6页：竞争格局特点
add_content_slide(prs, "竞争格局特点", [
    "三大竞争梯队",
    "",
    "第一梯队：互联网巨头",
    "• 百度、阿里、腾讯、字节凭借技术实力和生态优势占据主导地位",
    "• 阿里通义千问开源模型下载量超3亿次，衍生模型数超10万",
    "• 字节豆包日均Token调用量达12.7万亿次，市场份额46.4%",
    "",
    "第二梯队：AI原生企业",
    "• DeepSeek、月之暗面、智谱AI等快速崛起",
    "• DeepSeek以低成本高性能（训练成本仅GPT-4的1/27）颠覆行业",
    "• 在具身智能、AI基础设施等领域形成差异化竞争",
    "",
    "第三梯队：垂直领域企业",
    "• 商汤科技、科大讯飞等在特定领域展现优势",
    "• 聚焦行业解决方案，深耕垂直场景"
])

# 第7页：关键技术趋势
add_content_slide(prs, "三、关键技术趋势", [
    "大模型技术演进",
    "• 参数规模持续扩大：从百亿到万亿参数",
    "• MoE（混合专家）架构成为主流，动态分配计算资源",
    "• 推理能力显著提升：DeepSeek-R1等推理优化模型涌现",
    "• 成本持续下降：阿里通义千问主力模型降价97%",
    "",
    "多模态融合",
    "• 文本、图像、语音、视频全模态能力整合",
    "• 多模态AI被视为2024年重要趋势之一",
    "• 应用场景拓展至3D生成、智能客服等领域",
    "",
    "数据来源：中国信通院、量子位"
])

# 第8页：技术趋势（续）
add_content_slide(prs, "关键技术趋势（续）", [
    "工具生态完善",
    "• 分布式AI框架：DeepSpeed、Megatron、Colossal-AI",
    "• LLMOps平台：面向大模型全生命周期管理",
    "• 一体机产品：软硬件集成，降低应用门槛",
    "",
    "应用落地加速",
    "• 从百模大战向应用为王转变",
    "• Agent智能体成为新热点",
    "• 端侧AI寻求规模效应突破",
    "• 知识管理、营销、办公、教育场景应用最广",
    "",
    "开源生态繁荣",
    "• 阿里已开源200多款模型，成为全球第一AI开源模型",
    "• 开源成为技术突围的关键路径"
])

# 第9页：投资机会分析
add_content_slide(prs, "四、投资机会分析", [
    "融资市场概况（2024年）",
    "• AI行业融资总额超1000亿元",
    "• 上半年完成142次融资事件，已披露金额142.8亿元",
    "• 国家人工智能产业投资基金成立，出资600.6亿元",
    "",
    "热门投资方向",
    "• 基础层：AI芯片（11起融资）、算力基础设施",
    "• 模型层：通用大模型持续获得资本青睐",
    "• 应用层：文娱、知识管理、营销自动化",
    "",
    "典型案例",
    "• 月之暗面：2024年2月获投超10亿美金",
    "• DeepSeek：15天日活突破1500万，成为现象级应用",
    "",
    "数据来源：IT桔子、睿兽分析"
])

# 第10页：投资机会（续）
add_content_slide(prs, "重点投资赛道", [
    "高潜力赛道",
    "",
    "1. AI基础设施",
    "   • 智能算力需求持续增长，推理侧算力需求大幅上涨",
    "   • AI云市场2025年规模达223亿元，同比翻倍",
    "   • 预计2030年AI云市场达1930亿元",
    "",
    "2. 垂直行业应用",
    "   • 政务、医疗、能源等场景深度应用",
    "   • 企业级AI应用订阅制模式逐步成熟",
    "",
    "3. AI工具链",
    "   • LLMOps平台、一体机解决方案",
    "   • 降低企业AI应用门槛的工具产品",
    "",
    "4. 出海市场",
    "   • 东南亚和北美最受青睐，占比67.1%和62.1%"
])

# 第11页：总结与展望
add_content_slide(prs, "总结与展望", [
    "核心观点",
    "",
    "市场规模",
    "  2024年中国AI产业规模突破7000亿元，保持高速增长态势",
    "",
    "竞争格局",
    "  形成互联网巨头、AI原生企业、垂直领域企业三大梯队",
    "  阿里、字节、华为、腾讯、百度占据AI云市场主导地位",
    "",
    "技术趋势",
    "  大模型持续迭代，多模态融合、开源生态、应用落地加速",
    "",
    "投资机会",
    "  AI基础设施、垂直应用、工具链、出海市场具有较大潜力",
    "  2024年融资超1000亿元，资本持续看好AI赛道"
])

# 第12页：致谢页
slide_thanks = prs.slides.add_slide(prs.slide_layouts[6])
background = slide_thanks.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(240, 248, 255)
background.line.fill.background()

thanks_box = slide_thanks.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "谢谢观看"
thanks_para = thanks_frame.paragraphs[0]
thanks_para.font.size = Pt(48)
thanks_para.font.bold = True
thanks_para.font.color.rgb = COLOR_TITLE
thanks_para.alignment = PP_ALIGN.CENTER

# 保存文件
output_file = "2024年中国AI市场分析.pptx"
prs.save(output_file)
print(f"PPT已成功生成：{output_file}")
print(f"共生成 {len(prs.slides)} 页幻灯片")
