#!/usr/bin/env python3
"""
提取 PowerPoint 文件中所有形状（shape）的名称

使用方法:
    python extract_shapes.py template.pptx
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("❌ 需要安装 python-pptx: pip install python-pptx")
    sys.exit(1)


def extract_shapes(pptx_path: str):
    """提取 PPT 中所有 shape 的名称"""
    try:
        prs = Presentation(pptx_path)
        
        print(f"\n📄 文件: {pptx_path}")
        print(f"📊 总计 {len(prs.slides)} 页幻灯片\n")
        print("="*70)
        
        all_shapes = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes_in_slide = []
            
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "type": shape.shape_type,
                    "has_text": hasattr(shape, "text"),
                    "text_preview": ""
                }
                
                # 如果是文本框，获取文本预览
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    shape_info["text_preview"] = text[:50] + "..." if len(text) > 50 else text
                
                shapes_in_slide.append(shape_info)
                all_shapes.append({**shape_info, "slide": slide_idx})
            
            if shapes_in_slide:
                print(f"\n📍 幻灯片 {slide_idx}")
                print("-"*70)
                
                for shape_info in shapes_in_slide:
                    print(f"  • {shape_info['name']}")
                    if shape_info['text_preview']:
                        print(f"    └─ 内容: {shape_info['text_preview']}")
        
        # 统计信息
        print("\n" + "="*70)
        print(f"📊 统计信息:")
        print(f"  • 总形状数: {len(all_shapes)}")
        print(f"  • 包含文本的形状: {sum(1 for s in all_shapes if s['text_preview'])}")
        
        # 可编辑的 shape（有文本的）
        editable_shapes = [s for s in all_shapes if s['text_preview']]
        
        if editable_shapes:
            print(f"\n📝 可编辑的形状名称列表:")
            print("-"*70)
            for shape in editable_shapes:
                print(f"  • {shape['name']} (幻灯片 {shape['slide']})")
            
            # 生成示例配置
            print(f"\n💡 示例配置:")
            print("-"*70)
            print('```json')
            print('{')
            print('  "replacements": [')
            for i, shape in enumerate(editable_shapes[:5]):  # 只显示前5个
                comma = "," if i < min(len(editable_shapes), 5) - 1 else ""
                print(f'    {{"shape_name": "{shape["name"]}", "content": "新内容"}}{comma}')
            if len(editable_shapes) > 5:
                print(f'    // ... 还有 {len(editable_shapes) - 5} 个形状')
            print('  ]')
            print('}')
            print('```')
        
        print("\n" + "="*70)
        print("✅ 提取完成!\n")
        
    except FileNotFoundError:
        print(f"❌ 文件不存在: {pptx_path}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ 错误: {e}")
        sys.exit(1)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python extract_shapes.py <pptx_file>")
        print("示例: python extract_shapes.py template.pptx")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    extract_shapes(pptx_path)

