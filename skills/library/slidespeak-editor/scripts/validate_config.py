#!/usr/bin/env python3
"""
验证编辑配置是否有效

检查:
1. config JSON 格式是否正确
2. 所有 shape_name 是否存在于模板中
3. 内容长度是否合理

使用方法:
    python validate_config.py --template template.pptx --config config.json
"""

import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("❌ 需要安装 python-pptx: pip install python-pptx")
    sys.exit(1)


def get_all_shape_names(pptx_path: str):
    """获取 PPT 中所有 shape 名称"""
    prs = Presentation(pptx_path)
    shape_names = set()
    
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_names.add(shape.name)
    
    return shape_names


def validate_config(template_path: str, config_path: str):
    """验证配置"""
    print(f"\n{'='*70}")
    print(f"🔍 验证编辑配置")
    print(f"{'='*70}\n")
    
    # 1. 读取配置
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
    except FileNotFoundError:
        print(f"❌ 配置文件不存在: {config_path}")
        return False
    except json.JSONDecodeError as e:
        print(f"❌ 配置文件 JSON 格式错误: {e}")
        return False
    
    print(f"✅ 配置文件格式正确")
    
    # 2. 验证必需字段
    if "replacements" not in config:
        print(f"❌ 缺少必需字段: 'replacements'")
        return False
    
    replacements = config["replacements"]
    if not isinstance(replacements, list):
        print(f"❌ 'replacements' 必须是数组")
        return False
    
    if len(replacements) == 0:
        print(f"❌ 'replacements' 不能为空")
        return False
    
    print(f"✅ 包含 {len(replacements)} 个替换项")
    
    # 3. 验证每个替换项
    issues = []
    
    for i, replacement in enumerate(replacements, 1):
        if not isinstance(replacement, dict):
            issues.append(f"  • 替换项 {i}: 必须是对象")
            continue
        
        if "shape_name" not in replacement:
            issues.append(f"  • 替换项 {i}: 缺少 'shape_name' 字段")
            continue
        
        if "content" not in replacement:
            issues.append(f"  • 替换项 {i}: 缺少 'content' 字段")
            continue
        
        # 检查内容长度
        content = replacement["content"]
        if len(content) > 1000:
            issues.append(f"  • 替换项 {i} ({replacement['shape_name']}): 内容过长 ({len(content)} 字符)")
    
    if issues:
        print(f"\n⚠️  发现 {len(issues)} 个配置问题:")
        for issue in issues:
            print(issue)
        return False
    
    print(f"✅ 所有替换项格式正确")
    
    # 4. 验证 shape 名称是否存在于模板中
    try:
        template_shapes = get_all_shape_names(template_path)
        print(f"✅ 模板包含 {len(template_shapes)} 个形状")
        
        missing_shapes = []
        for replacement in replacements:
            shape_name = replacement["shape_name"]
            if shape_name not in template_shapes:
                missing_shapes.append(shape_name)
        
        if missing_shapes:
            print(f"\n❌ 以下 {len(missing_shapes)} 个形状在模板中不存在:")
            for shape_name in missing_shapes:
                print(f"  • {shape_name}")
            print(f"\n💡 提示: 使用 extract_shapes.py 查看模板中的所有形状名称")
            return False
        
        print(f"✅ 所有形状名称都存在于模板中")
        
    except FileNotFoundError:
        print(f"❌ 模板文件不存在: {template_path}")
        return False
    except Exception as e:
        print(f"❌ 读取模板时出错: {e}")
        return False
    
    # 5. 显示验证摘要
    print(f"\n{'='*70}")
    print(f"📊 验证摘要:")
    print(f"  • 替换数量: {len(replacements)}")
    print(f"  • 平均内容长度: {sum(len(r['content']) for r in replacements) // len(replacements)} 字符")
    print(f"  • 状态: ✅ 验证通过")
    print(f"{'='*70}\n")
    
    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="验证 SlideSpeak 编辑配置")
    parser.add_argument("--template", required=True, help="模板 PPT 文件路径")
    parser.add_argument("--config", required=True, help="配置 JSON 文件路径")
    
    args = parser.parse_args()
    
    success = validate_config(args.template, args.config)
    sys.exit(0 if success else 1)

