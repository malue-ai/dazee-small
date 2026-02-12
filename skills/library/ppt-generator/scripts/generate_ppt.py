"""
PPT 生成脚本 - 使用 python-pptx 创建专业演示文稿

这个脚本在 Claude Code Execution 环境中运行，
Claude 会根据用户需求调用这个脚本生成 PPT。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import json
from typing import Dict, List, Any, Optional
from dataclasses import dataclass
from pathlib import Path


# ============================================================
# 主题配置
# ============================================================
THEMES = {
    "business_blue": {
        "primary": RgbColor(0x00, 0x66, 0xB3),      # 深蓝
        "secondary": RgbColor(0x00, 0x99, 0xE5),    # 亮蓝
        "accent": RgbColor(0xFF, 0x99, 0x00),       # 橙色点缀
        "background": RgbColor(0xFF, 0xFF, 0xFF),   # 白色背景
        "text_dark": RgbColor(0x33, 0x33, 0x33),    # 深灰文字
        "text_light": RgbColor(0xFF, 0xFF, 0xFF),   # 白色文字
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
    },
    "tech_dark": {
        "primary": RgbColor(0x1A, 0x1A, 0x2E),      # 深紫黑
        "secondary": RgbColor(0x16, 0x21, 0x3E),    # 深蓝
        "accent": RgbColor(0x0F, 0xF0, 0xFC),       # 青色
        "background": RgbColor(0x1A, 0x1A, 0x2E),   # 深色背景
        "text_dark": RgbColor(0xE9, 0xE9, 0xE9),    # 浅灰文字
        "text_light": RgbColor(0xFF, 0xFF, 0xFF),   # 白色文字
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
    },
    "nature_green": {
        "primary": RgbColor(0x2E, 0x7D, 0x32),      # 深绿
        "secondary": RgbColor(0x4C, 0xAF, 0x50),    # 中绿
        "accent": RgbColor(0xFF, 0xC1, 0x07),       # 金色
        "background": RgbColor(0xF5, 0xF5, 0xF5),   # 浅灰背景
        "text_dark": RgbColor(0x21, 0x21, 0x21),    # 深色文字
        "text_light": RgbColor(0xFF, 0xFF, 0xFF),   # 白色文字
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
    },
    "elegant_gray": {
        "primary": RgbColor(0x42, 0x42, 0x42),      # 深灰
        "secondary": RgbColor(0x75, 0x75, 0x75),    # 中灰
        "accent": RgbColor(0xC6, 0x28, 0x28),       # 红色点缀
        "background": RgbColor(0xFA, 0xFA, 0xFA),   # 极浅灰
        "text_dark": RgbColor(0x21, 0x21, 0x21),    # 深色文字
        "text_light": RgbColor(0xFF, 0xFF, 0xFF),   # 白色文字
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
    },
    "vibrant_orange": {
        "primary": RgbColor(0xFF, 0x57, 0x22),      # 深橙
        "secondary": RgbColor(0xFF, 0x98, 0x00),    # 橙色
        "accent": RgbColor(0x00, 0x96, 0x88),       # 青绿
        "background": RgbColor(0xFF, 0xFF, 0xFF),   # 白色背景
        "text_dark": RgbColor(0x33, 0x33, 0x33),    # 深色文字
        "text_light": RgbColor(0xFF, 0xFF, 0xFF),   # 白色文字
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
    },
}


@dataclass
class SlideContent:
    """幻灯片内容定义"""
    slide_type: str  # title, content, chart, image, section, ending
    title: str
    subtitle: Optional[str] = None
    body: Optional[List[str]] = None  # 项目符号列表
    chart_data: Optional[Dict[str, Any]] = None
    image_url: Optional[str] = None
    notes: Optional[str] = None


class PPTGenerator:
    """PPT 生成器"""
    
    def __init__(self, theme: str = "business_blue"):
        """
        初始化 PPT 生成器
        
        Args:
            theme: 主题名称
        """
        self.prs = Presentation()
        # 设置 16:9 宽屏比例
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        self.theme = THEMES.get(theme, THEMES["business_blue"])
        self.theme_name = theme
    
    def add_title_slide(self, title: str, subtitle: str = "", author: str = "", company: str = ""):
        """
        添加标题幻灯片
        
        Args:
            title: 主标题
            subtitle: 副标题
            author: 作者
            company: 公司
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景色块（上半部分）
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["primary"]
        shape.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        title_para.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.2),
                Inches(12.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.theme["text_light"]
            subtitle_para.font.name = self.theme["body_font"]
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # 作者和公司信息
        if author or company:
            info_text = f"{author}  |  {company}" if author and company else (author or company)
            info_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5),
                Inches(12.333), Inches(0.5)
            )
            info_frame = info_box.text_frame
            info_para = info_frame.paragraphs[0]
            info_para.text = info_text
            info_para.font.size = Pt(16)
            info_para.font.color.rgb = self.theme["text_dark"]
            info_para.font.name = self.theme["body_font"]
            info_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_section_slide(self, title: str, subtitle: str = ""):
        """
        添加章节分隔幻灯片
        
        Args:
            title: 章节标题
            subtitle: 章节副标题
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 全页背景色
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["secondary"]
        shape.line.fill.background()
        
        # 章节编号装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(3.2),
            Inches(2), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()
        
        # 章节标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5),
            Inches(11), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        
        if subtitle:
            p = title_frame.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = self.theme["text_light"]
            p.font.name = self.theme["body_font"]
            p.space_before = Pt(12)
        
        return slide
    
    def add_content_slide(self, title: str, bullet_points: List[str], subtitle: str = ""):
        """
        添加内容幻灯片（项目符号列表）
        
        Args:
            title: 幻灯片标题
            bullet_points: 项目符号列表
            subtitle: 副标题（可选）
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary"]
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        
        # 内容区域
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8),
            Inches(11.5), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.theme["text_dark"]
            p.font.name = self.theme["body_font"]
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        
        return slide
    
    def add_chart_slide(self, title: str, chart_data: Dict[str, Any], chart_type: str = "column"):
        """
        添加图表幻灯片
        
        Args:
            title: 幻灯片标题
            chart_data: 图表数据 {"categories": [...], "series": [{"name": ..., "values": [...]}]}
            chart_type: 图表类型 (column, bar, line, pie)
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary"]
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        
        # 创建图表数据
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get("categories", [])
        
        for series in chart_data.get("series", []):
            chart_data_obj.add_series(series["name"], series["values"])
        
        # 确定图表类型
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # 添加图表
        x, y, cx, cy = Inches(1), Inches(1.8), Inches(11), Inches(5)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data_obj
        ).chart
        
        return slide
    
    def add_two_column_slide(self, title: str, left_content: List[str], right_content: List[str]):
        """
        添加双栏内容幻灯片
        
        Args:
            title: 幻灯片标题
            left_content: 左侧内容
            right_content: 右侧内容
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary"]
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        
        # 左侧内容
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8),
            Inches(5.8), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, point in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme["text_dark"]
            p.font.name = self.theme["body_font"]
            p.space_before = Pt(10)
        
        # 右侧内容
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.8),
            Inches(5.8), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, point in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme["text_dark"]
            p.font.name = self.theme["body_font"]
            p.space_before = Pt(10)
        
        return slide
    
    def add_ending_slide(self, title: str = "谢谢观看", subtitle: str = "", contact: str = ""):
        """
        添加结束幻灯片
        
        Args:
            title: 结束语
            subtitle: 副标题
            contact: 联系方式
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 全页背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["primary"]
        shape.line.fill.background()
        
        # 结束语
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme["text_light"]
        title_para.font.name = self.theme["title_font"]
        title_para.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            p = title_frame.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.theme["text_light"]
            p.font.name = self.theme["body_font"]
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(20)
        
        # 联系方式
        if contact:
            contact_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5),
                Inches(12.333), Inches(0.5)
            )
            contact_frame = contact_box.text_frame
            contact_para = contact_frame.paragraphs[0]
            contact_para.text = contact
            contact_para.font.size = Pt(14)
            contact_para.font.color.rgb = self.theme["text_light"]
            contact_para.font.name = self.theme["body_font"]
            contact_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def save(self, filename: str) -> str:
        """
        保存 PPT 文件
        
        Args:
            filename: 文件名（不含扩展名）
            
        Returns:
            保存的文件路径
        """
        if not filename.endswith('.pptx'):
            filename = f"{filename}.pptx"
        
        self.prs.save(filename)
        return filename


def create_ppt_from_config(config: Dict[str, Any]) -> str:
    """
    从配置创建 PPT
    
    Args:
        config: PPT 配置，包含：
            - title: 演示文稿标题
            - theme: 主题名称
            - author: 作者
            - company: 公司
            - slides: 幻灯片配置列表
            
    Returns:
        生成的文件路径
    """
    theme = config.get("theme", "business_blue")
    generator = PPTGenerator(theme=theme)
    
    slides = config.get("slides", [])
    
    for i, slide_config in enumerate(slides):
        slide_type = slide_config.get("type", "content")
        
        if slide_type == "title":
            generator.add_title_slide(
                title=slide_config.get("title", "演示文稿"),
                subtitle=slide_config.get("subtitle", ""),
                author=config.get("author", ""),
                company=config.get("company", "")
            )
        
        elif slide_type == "section":
            generator.add_section_slide(
                title=slide_config.get("title", f"第 {i+1} 部分"),
                subtitle=slide_config.get("subtitle", "")
            )
        
        elif slide_type == "content":
            generator.add_content_slide(
                title=slide_config.get("title", "内容"),
                bullet_points=slide_config.get("bullets", []),
                subtitle=slide_config.get("subtitle", "")
            )
        
        elif slide_type == "chart":
            generator.add_chart_slide(
                title=slide_config.get("title", "图表"),
                chart_data=slide_config.get("chart_data", {}),
                chart_type=slide_config.get("chart_type", "column")
            )
        
        elif slide_type == "two_column":
            generator.add_two_column_slide(
                title=slide_config.get("title", "对比"),
                left_content=slide_config.get("left", []),
                right_content=slide_config.get("right", [])
            )
        
        elif slide_type == "ending":
            generator.add_ending_slide(
                title=slide_config.get("title", "谢谢观看"),
                subtitle=slide_config.get("subtitle", ""),
                contact=slide_config.get("contact", "")
            )
    
    # 保存文件
    filename = config.get("filename", config.get("title", "presentation"))
    output_path = generator.save(filename)
    
    return output_path


# ============================================================
# 示例用法（Claude 会直接调用这些函数）
# ============================================================
if __name__ == "__main__":
    # 示例配置
    sample_config = {
        "title": "2024 Q4 销售报告",
        "theme": "business_blue",
        "author": "张三",
        "company": "示例公司",
        "filename": "Q4_Sales_Report",
        "slides": [
            {
                "type": "title",
                "title": "2024 Q4 销售报告",
                "subtitle": "年度业绩回顾与展望"
            },
            {
                "type": "section",
                "title": "01 季度概览",
                "subtitle": "关键数据与里程碑"
            },
            {
                "type": "content",
                "title": "季度亮点",
                "bullets": [
                    "总销售额达到 2100 万元，同比增长 40%",
                    "新客户数量突破 500 家",
                    "客户满意度提升至 95%",
                    "成功进入 3 个新市场区域"
                ]
            },
            {
                "type": "chart",
                "title": "季度销售趋势",
                "chart_type": "column",
                "chart_data": {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "销售额(万元)", "values": [1200, 1500, 1800, 2100]}
                    ]
                }
            },
            {
                "type": "two_column",
                "title": "机遇与挑战",
                "left": [
                    "市场需求持续增长",
                    "新产品线反响良好",
                    "团队能力显著提升"
                ],
                "right": [
                    "竞争对手加大投入",
                    "供应链成本上升",
                    "人才招聘仍有缺口"
                ]
            },
            {
                "type": "content",
                "title": "下季度计划",
                "bullets": [
                    "推出 2 款新产品",
                    "扩展线上销售渠道",
                    "加强客户关系管理",
                    "优化供应链效率"
                ]
            },
            {
                "type": "ending",
                "title": "谢谢观看",
                "subtitle": "期待与您共创辉煌",
                "contact": "contact@example.com"
            }
        ]
    }
    
    # 生成 PPT
    output_file = create_ppt_from_config(sample_config)
    print(f"✅ PPT 已生成: {output_file}")


